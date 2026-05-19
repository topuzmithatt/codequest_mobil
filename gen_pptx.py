from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(15, 23, 42)
ACCENT = RGBColor(99, 102, 241)
ACCENT2 = RGBColor(16, 185, 129)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)
CARD_BG = RGBColor(30, 41, 59)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_shape(slide, left, top, w, h, color=CARD_BG, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

# ============ SLIDE 1 - TITLE ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_shape(sl, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(20,27,50))
add_text(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2), "🏆 CodeQuest", 52, ACCENT, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3), Inches(11), Inches(0.8), "Oyunlaştırılmış Yazılım Eğitim Platformu", 32, WHITE, True, PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(4.2), Inches(9), Inches(0.6), "Next.js 16 · React 19 · Prisma 7 · Supabase · Groq LLaMA 3.3 · Monaco Editor", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(sl, Inches(2), Inches(5.5), Inches(9), Inches(0.6), "Bitirme Projesi Sunumu", 20, ACCENT2, True, PP_ALIGN.CENTER)

# ============ SLIDE 2 - PROJE OZETI ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "🎯 Proje Özeti", 36, ACCENT, True)
add_shape(sl, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5), CARD_BG)
add_text(sl, Inches(1.2), Inches(1.7), Inches(10.5), Inches(2.2),
    "CodeQuest, kullanıcıların Python, JavaScript, Java ve SQL dillerini interaktif görevler, "
    "yapay zeka destekli kişiselleştirme ve oyunlaştırma mekanikleriyle öğrenmesini sağlayan "
    "modern bir web tabanlı eğitim platformudur.\n\n"
    "Her kullanıcıya onboarding sürecinde belirlenen hedef, deneyim seviyesi ve seçilen dile göre "
    "kişiselleştirilmiş bir öğrenme yolu oluşturulur.", 16, WHITE)

add_text(sl, Inches(0.8), Inches(4.3), Inches(11), Inches(0.6), "Çözülen Problemler", 24, ACCENT2, True)
problems = [
    ("Pasif öğrenme → düşük motivasyon", "XP, seviye, rozet ve seri mekanikleriyle aktif katılım"),
    ("Herkese aynı müfredat", "AI destekli kişiselleştirilmiş öğrenme yolu"),
    ("Kod çalıştırma ortamı kurma zorluğu", "Tarayıcı içi sandbox kod çalıştırma"),
    ("Geri bildirim eksikliği", "AI Code Reviewer ile 5 eksenli otomatik değerlendirme"),
]
for i, (prob, sol) in enumerate(problems):
    y = 5.0 + i * 0.55
    add_shape(sl, Inches(0.8), Inches(y), Inches(5.5), Inches(0.45), RGBColor(40,20,20))
    add_text(sl, Inches(1.0), Inches(y+0.05), Inches(5.2), Inches(0.35), f"❌ {prob}", 13, RGBColor(248,113,113))
    add_shape(sl, Inches(6.5), Inches(y), Inches(5.8), Inches(0.45), RGBColor(20,40,30))
    add_text(sl, Inches(6.7), Inches(y+0.05), Inches(5.5), Inches(0.35), f"✅ {sol}", 13, ACCENT2)

# ============ SLIDE 3 - OZELLIKLER ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "✨ Temel Özellikler", 36, ACCENT, True)

features = [
    ("🎮 Oyunlaştırma", ["XP & Seviye Sistemi (10 kademe)", "❤️ Can Sistemi (5 hak, saatte yenilenir)", "🔥 Günlük Seri (Streak × 1.5 XP)", "🏅 Rozet Sistemi", "📊 Haftalık XP Takibi"]),
    ("🤖 Yapay Zeka", ["AI Onboarding - kişiselleştirilmiş yol", "3 Aşamalı Görev Üretimi", "AI Code Reviewer (5 eksen)", "Akıllı İpucu Sistemi"]),
    ("💻 Editör & Terminal", ["Monaco Editor (IntelliSense)", "xterm.js Terminal", "Görev Paneli & Test Sonuçları", "Radar Grafik Değerlendirme"]),
    ("📁 Portföy & Profil", ["XP, seviye, rozet takibi", "Herkese açık portföy", "Çoklu dil desteği"]),
]
for i, (title, items) in enumerate(features):
    x = 0.8 + i * 3.1
    add_shape(sl, Inches(x), Inches(1.5), Inches(2.9), Inches(5.2), CARD_BG)
    add_text(sl, Inches(x+0.2), Inches(1.6), Inches(2.5), Inches(0.5), title, 18, ACCENT2, True)
    for j, item in enumerate(items):
        add_text(sl, Inches(x+0.2), Inches(2.3+j*0.55), Inches(2.5), Inches(0.5), f"• {item}", 12, WHITE)

# ============ SLIDE 4 - TECH STACK ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "🛠️ Tech Stack", 36, ACCENT, True)

# Frontend
add_text(sl, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5), "Frontend", 22, ACCENT2, True)
fe_items = [("Next.js 16","App Router, SSR"),("React 19","Bileşen tabanlı UI"),("TypeScript 5","Tip güvenliği"),
    ("Tailwind CSS 4","Utility-first"),("Monaco Editor","Kod editörü"),("xterm.js","Terminal"),("Socket.IO Client","WebSocket")]
for i,(t,d) in enumerate(fe_items):
    add_shape(sl, Inches(0.8), Inches(1.9+i*0.65), Inches(5.5), Inches(0.55), CARD_BG)
    add_text(sl, Inches(1.0), Inches(1.95+i*0.65), Inches(2.2), Inches(0.45), t, 14, WHITE, True)
    add_text(sl, Inches(3.3), Inches(1.95+i*0.65), Inches(2.8), Inches(0.45), d, 13, GRAY)

# Backend
add_text(sl, Inches(7), Inches(1.3), Inches(5), Inches(0.5), "Backend & Veritabanı", 22, ACCENT2, True)
be_items = [("Next.js API Routes","RESTful API"),("Express.js 5","Execution Server"),("Socket.IO 4.8","Gerçek zamanlı"),
    ("Groq SDK","LLaMA 3.3 70B"),("Prisma ORM 7","DB erişim katmanı"),("Supabase","PostgreSQL + Auth"),("child_process","Sandbox çalıştırma")]
for i,(t,d) in enumerate(be_items):
    add_shape(sl, Inches(7), Inches(1.9+i*0.65), Inches(5.5), Inches(0.55), CARD_BG)
    add_text(sl, Inches(7.2), Inches(1.95+i*0.65), Inches(2.2), Inches(0.45), t, 14, WHITE, True)
    add_text(sl, Inches(9.5), Inches(1.95+i*0.65), Inches(2.8), Inches(0.45), d, 13, GRAY)

# ============ SLIDE 5 - MIMARI ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "⚙️ Sistem Mimarisi", 36, ACCENT, True)

boxes = [
    (1, 0.5, 3.5, 2, "🖥️ Kullanıcı (Tarayıcı)", ["Monaco Editor", "xterm.js Terminal", "Challenge Panel"], RGBColor(30,58,95)),
    (5.5, 0.5, 3.5, 2, "🚀 Execution Server :3001", ["child_process.spawn()", "Python / Node / Java / SQLite", ".temp_exec/ geçici dosyalar"], RGBColor(55,30,30)),
    (1, 3.5, 3.5, 2.5, "📡 Next.js API :3000", ["/api/submit → Test", "/api/run → Çalıştır", "/api/onboarding → AI Yolu", "/api/ai-hint → İpucu", "/api/hearts → Can"], RGBColor(30,50,40)),
    (5.5, 3.5, 3.5, 1.5, "🗄️ Prisma ORM", ["PostgreSQL (Supabase)", "12 Model, 5 Enum"], RGBColor(35,35,55)),
    (5.5, 5.2, 3.5, 1.3, "🤖 Groq API", ["LLaMA 3.3 70B", "Görev üretimi & Değerlendirme"], RGBColor(50,30,45)),
]
for (x,y,w,h,title,items,color) in boxes:
    add_shape(sl, Inches(x), Inches(y+1), Inches(w), Inches(h), color)
    add_text(sl, Inches(x+0.15), Inches(y+1.05), Inches(w-0.3), Inches(0.4), title, 14, ACCENT2, True)
    for j, item in enumerate(items):
        add_text(sl, Inches(x+0.15), Inches(y+1.5+j*0.35), Inches(w-0.3), Inches(0.3), f"• {item}", 11, WHITE)

add_text(sl, Inches(10), Inches(1.5), Inches(2.5), Inches(5), "WebSocket ↔️\n\nHTTP REST ↔️\n\nPrisma Client ↔️", 14, GRAY, False, PP_ALIGN.CENTER)

# ============ SLIDE 6 - SANDBOX AKISI ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "🔄 Kod Çalıştırma Akışı (Sandbox)", 36, ACCENT, True)

steps = [
    "1. Kullanıcı Monaco Editor'de kod yazar",
    "2. Socket.IO ile 'runCode' eventi → Execution Server",
    "3. Benzersiz sessionId üretilir (crypto.randomBytes)",
    "4. Kod geçici dosyaya yazılır (.temp_exec/)",
    "5. Dile göre komut belirlenir (python/node/javac/sqlite3)",
    "6. child_process.spawn() ile süreç başlatılır",
    "7. stdout/stderr → Socket.IO → Terminal (gerçek zamanlı)",
    "8. Süreç tamamlanınca 'finished' eventi gönderilir",
    "9. Geçici dosyalar temizlenir | Timeout: 10s → SIGKILL",
]
for i, step in enumerate(steps):
    y = 1.4 + i * 0.62
    c = ACCENT if i % 2 == 0 else ACCENT2
    add_shape(sl, Inches(1), Inches(y), Inches(11), Inches(0.52), CARD_BG)
    add_text(sl, Inches(1.3), Inches(y+0.05), Inches(10.5), Inches(0.42), step, 15, c, i==0)

# ============ SLIDE 7 - GOREV GONDERIMI ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "📝 Görev Gönderimi ve Doğrulama", 36, ACCENT, True)

add_shape(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), CARD_BG)
add_text(sl, Inches(1), Inches(1.6), Inches(5), Inches(0.4), "Test Akışı", 20, ACCENT2, True)
test_steps = [
    "POST /api/submit → Kod + challengeId",
    "pistonClient.runTests() çağrılır",
    "Her test case için kod ayrı çalıştırılır",
    "stdin → test input verilir",
    "stdout ↔ expectedOutput karşılaştırılır",
    "Normalize: trim, lowercase, Türkçe karakter",
    "Float tolerans: %0.001",
    "Sonuç: { passedCount, totalCount, allPassed }",
]
for i, s in enumerate(test_steps):
    add_text(sl, Inches(1.2), Inches(2.2+i*0.55), Inches(4.8), Inches(0.45), f"→ {s}", 13, WHITE)

add_shape(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(5.5), CARD_BG)
add_text(sl, Inches(7.2), Inches(1.6), Inches(5), Inches(0.4), "Gamification Engine", 20, ACCENT2, True)
gam_steps = [
    "✅ allPassed = true:",
    "   • awardXP() → XP kazandır",
    "   • Seviye kontrolü, rozet kontrolü",
    "   • updateStreak() → Seri güncelle",
    "   • XP çarpanı ayarla (×1.5)",
    "❌ allPassed = false:",
    "   • loseHeart() → Can düşür",
    "   • Submission status → FAILED",
    "🤖 AI Code Reviewer:",
    "   • 5 eksenli değerlendirme (Groq)",
]
for i, s in enumerate(gam_steps):
    c = ACCENT2 if "✅" in s else (RGBColor(248,113,113) if "❌" in s else (RGBColor(129,140,248) if "🤖" in s else WHITE))
    add_text(sl, Inches(7.4), Inches(2.2+i*0.5), Inches(4.8), Inches(0.4), s, 13, c)

# ============ SLIDE 8 - AI PIPELINE ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "🤖 AI Görev Üretimi (3 Aşamalı Pipeline)", 36, ACCENT, True)

stages = [
    ("Aşama 1 — Yaratıcı Üretim", "temperature: 0.7", "Groq → Soru metni, starterCode, test case input'ları, ipuçları", RGBColor(99,102,241)),
    ("Aşama 2 — Referans Çözüm", "temperature: 0.1", "Groq → solutionCode (tam çalışır çözüm)", RGBColor(16,185,129)),
    ("Aşama 3 — Doğrulama", "runtime execution", "solutionCode her test input'u ile çalıştırılır\n✅ Başarılı → stdout expectedOutput olarak kaydedilir\n❌ Başarısız → Fallback: Groq'a expectedOutput sorulur", RGBColor(251,191,36)),
]
for i, (title, sub, desc, color) in enumerate(stages):
    y = 1.5 + i * 1.9
    add_shape(sl, Inches(1), Inches(y), Inches(11), Inches(1.7), CARD_BG)
    add_text(sl, Inches(1.3), Inches(y+0.1), Inches(6), Inches(0.5), title, 22, color, True)
    add_text(sl, Inches(8), Inches(y+0.1), Inches(3.5), Inches(0.4), sub, 14, GRAY, False, PP_ALIGN.RIGHT)
    add_text(sl, Inches(1.3), Inches(y+0.7), Inches(10.2), Inches(0.9), desc, 14, WHITE)

# ============ SLIDE 9 - VERITABANI ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "🗄️ Veritabanı Şeması", 36, ACCENT, True)
add_text(sl, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4), "12 Model · 5 Enum · PostgreSQL (Supabase) · Prisma ORM", 16, GRAY)

models = [
    ("User", "hearts, xp, level", RGBColor(99,102,241)),
    ("LearningPath", "language, difficulty, topicsOrder", RGBColor(16,185,129)),
    ("Challenge", "testCases[], difficulty, xpReward", RGBColor(251,191,36)),
    ("Submission", "status, xpEarned, code", RGBColor(236,72,153)),
    ("ReviewResult", "correctness, readability, overallScore", RGBColor(139,92,246)),
    ("TestCase", "input, expectedOutput", RGBColor(14,165,233)),
    ("Streak", "current, longest", RGBColor(249,115,22)),
    ("Badge", "type, icon", RGBColor(34,197,94)),
]
for i, (name, fields, color) in enumerate(models):
    col = i % 4
    row = i // 4
    x = 0.8 + col * 3.1
    y = 2.0 + row * 2.5
    add_shape(sl, Inches(x), Inches(y), Inches(2.9), Inches(2.0), CARD_BG)
    add_text(sl, Inches(x+0.2), Inches(y+0.1), Inches(2.5), Inches(0.5), name, 18, color, True)
    add_text(sl, Inches(x+0.2), Inches(y+0.7), Inches(2.5), Inches(1.1), fields, 13, GRAY)

# ============ SLIDE 10 - DILLER ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "💻 Desteklenen Programlama Dilleri", 36, ACCENT, True)

langs = [
    ("🐍 Python", "python -u", "stdin/stdout, UTF-8 desteği", RGBColor(55,118,171)),
    ("📜 JavaScript", "node", "ES6+, readline desteği", RGBColor(247,223,30)),
    ("☕ Java", "javac + java", "Derleme + çalıştırma, Scanner", RGBColor(237,28,36)),
    ("🗃️ SQL", "Python sqlite3", "In-memory DB, çoklu sorgu", RGBColor(0,150,136)),
]
for i, (name, runtime, features, color) in enumerate(langs):
    x = 0.8 + i * 3.1
    add_shape(sl, Inches(x), Inches(1.5), Inches(2.9), Inches(4.5), CARD_BG)
    add_text(sl, Inches(x+0.2), Inches(1.7), Inches(2.5), Inches(0.6), name, 24, color, True, PP_ALIGN.CENTER)
    add_text(sl, Inches(x+0.2), Inches(2.5), Inches(2.5), Inches(0.4), "Çalıştırma:", 13, GRAY)
    add_text(sl, Inches(x+0.2), Inches(2.9), Inches(2.5), Inches(0.5), runtime, 16, WHITE, True)
    add_text(sl, Inches(x+0.2), Inches(3.6), Inches(2.5), Inches(0.4), "Özellikler:", 13, GRAY)
    add_text(sl, Inches(x+0.2), Inches(4.0), Inches(2.5), Inches(1.5), features, 14, WHITE)

# ============ SLIDE 11 - KURULUM ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_text(sl, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), "🚀 Kurulum ve Çalıştırma", 36, ACCENT, True)

setup_steps = [
    ("1", "git clone & cd codequest"),
    ("2", "npm install"),
    ("3", "cp .env.example .env → Değişkenleri yapılandır"),
    ("4", "npx prisma db push && npx prisma generate"),
    ("5", "npm run dev"),
    ("6", "http://localhost:3000 → Tarayıcıda açın"),
]
for i, (num, desc) in enumerate(setup_steps):
    y = 1.5 + i * 0.85
    add_shape(sl, Inches(1), Inches(y), Inches(0.7), Inches(0.65), ACCENT)
    add_text(sl, Inches(1), Inches(y+0.1), Inches(0.7), Inches(0.45), num, 22, WHITE, True, PP_ALIGN.CENTER)
    add_shape(sl, Inches(1.9), Inches(y), Inches(10), Inches(0.65), CARD_BG)
    add_text(sl, Inches(2.1), Inches(y+0.1), Inches(9.5), Inches(0.45), desc, 16, WHITE)

add_text(sl, Inches(1), Inches(6.8), Inches(11), Inches(0.5), "Ön Gereksinimler: Node.js 18+ · npm 9+ · Python 3.8+ · Java JDK 11+ (opsiyonel) · Git 2.0+", 14, GRAY)

# ============ SLIDE 12 - TESEKKUR ============
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
add_shape(sl, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(20,27,50))
add_text(sl, Inches(1), Inches(2), Inches(11), Inches(1), "Teşekkürler! 🎉", 52, ACCENT, True, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(3.5), Inches(11), Inches(0.8), "CodeQuest ile kodlamayı keşfetmeye başlayın! 🚀", 28, WHITE, False, PP_ALIGN.CENTER)
add_text(sl, Inches(1), Inches(5), Inches(11), Inches(0.6), "Sorularınız için hazırım.", 20, GRAY, False, PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(__file__), "CodeQuest_Sunum.pptx")
prs.save(out)
print(f"PPTX saved: {out}")
