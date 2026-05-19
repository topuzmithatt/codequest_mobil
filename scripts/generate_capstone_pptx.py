"""
CodeQuest bitirme sunumu — kabaca PPTX üretir.
Çalıştır: python scripts/generate_capstone_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent.parent / "presentations" / "CodeQuest_Bitirme_Sunumu.pptx"


def add_title_slide(prs, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullets(prs, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "CodeQuest",
        "İnteraktif Kodlama Öğrenme Platformu\nBitirme Projesi\n[Ad Soyad] · [Danışman] · [Üniversite / Bölüm] · 2026",
    )

    add_bullets(
        prs,
        "Problem",
        [
            "Yazılım öğrenmek isteyenler için yapılandırılmış, geri bildirim veren ortam ihtiyacı",
            "Motivasyon ve ilerleme takibinin (gamification) öğrenmeyi desteklemesi",
            "Farklı deneyim seviyelerine uyumlu içerik ve alıştırma zorluğu",
        ],
    )

    add_bullets(
        prs,
        "Amaç ve Kapsam",
        [
            "Web tabanlı bir platform: öğrenme yolları, kod editörü, testler",
            "Kullanıcı profili, gamification (XP, seviye, can, rozetler)",
            "Yapay zekâ destekli soru üretimi ve ipuçları (Groq)",
            "Kod çalıştırma: Piston API ve yerel yürütme sunucusu (Socket.io)",
        ],
    )

    add_bullets(
        prs,
        "Benzer Çalışmalar [Güncelle]",
        [
            "Codecademy, freeCodeCamp, LeetCode vb. — kısa karşılaştırma maddeleri ekleyin",
            "Bu projenin farkı: [tez cümleniz]",
        ],
    )

    add_bullets(
        prs,
        "Sistem Mimarisi",
        [
            "İstemci: Next.js (App Router), React, TypeScript",
            "Sunucu: Next.js API Routes, Express tabanlı exec-server",
            "Kimlik ve depolama: Supabase Auth + PostgreSQL",
            "ORM: Prisma 7 + bağlantı havuzu",
        ],
    )

    add_bullets(
        prs,
        "Öne Çıkan Özellikler",
        [
            "Onboarding ve dil / hedef seçimi",
            "Konu bazlı öğrenme ve challenge akışı",
            "Monaco editör, terminal entegrasyonu, test case sonuçları",
            "Profil, portföy ve gönderim görünürlüğü",
        ],
    )

    add_bullets(
        prs,
        "Veri Modeli (Özet)",
        [
            "Kullanıcı, öğrenme yolu, challenge, gönderim, gamification alanları",
            "Prisma şeması ile ilişkisel model — detay: teknik dokümantasyon / diyagram ekleyin",
        ],
    )

    add_bullets(
        prs,
        "Güvenlik ve Dağıtım",
        [
            "Oturum: Supabase ile sunucu tarafı doğrulama",
            "Ortam değişkenleri ile API anahtarları (client’ta expose edilmeyenler)",
            "exec-server: geliştirme odaklı; üretimde ayrı süreç ve kısıtlamalar — [tez notu]",
        ],
    )

    add_bullets(
        prs,
        "Test / Değerlendirme [Güncelle]",
        [
            "Fonksiyonel test senaryoları",
            "Kullanılabilirlik veya anket — varsa sonuç özeti",
        ],
    )

    add_bullets(
        prs,
        "Sonuç",
        [
            "Hedeflenen öğrenme akışı ve teknik bileşenler hayata geçirildi",
            "Gelecek çalışmalar: daha sıkı sandbox, mobil, ek diller, analitik",
        ],
    )

    add_title_slide(prs, "Teşekkürler", "Sorular")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Kaydedildi: {OUT}")


if __name__ == "__main__":
    main()
